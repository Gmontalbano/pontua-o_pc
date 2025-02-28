import collections
import collections.abc
from pptx import Presentation
from pptx import Presentation
from pptx.util import Pt, Cm
import aspose.slides as slides
from PIL import Image
import os
import pandas as pd

DEBUG = True
SLIDE_HEIGHT = 5716800
SLIDE_WIDTH = 9122400

GREY = (245, 245, 245)
DAKER_GREY = (200, 200, 200)
WHITE = (255, 255, 255)
BLACK = (0, 0, 0)


def replace_text(ppt, search_string, replace_string, how='all'):
    """
    Buil to change keyswords to data, in this case numbers
    Parameters
    ----------
    ppt : Presentation type
        slides
    test_string : Text
        key we need to search
    replace_string : Text
        data to put into the slide
    slide_n : int or None
        if has a number just replace on the slide with this index

    Returns
    -------
    pptx with the data in the places of the keys

    """
    assert (how in ['all', 'first'] or type(how) == int)

    if type(how) == int:
        slide_loop = [ppt.slides[how]]
    else:
        slide_loop = ppt.slides

    for slide in slide_loop:
        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape.text.find(search_string) != -1:
                    text_frame = shape.text_frame

                    for para in text_frame.paragraphs:
                        try:
                            cur_text = para.runs[0].text
                            # print(cur_text)

                            if str(search_string) in cur_text:

                                if DEBUG:
                                    print(cur_text, '->', replace_string)

                                new_text = cur_text.replace(
                                    str(search_string), str(replace_string))
                                para.runs[0].text = new_text

                                if how == 'first':
                                    return ppt

                        except Exception as e:
                            pass
    return ppt


def delete_shape(slide, shape):
    '''
    Function
    ----------
    lorem ipsum dolor sit amet

    Parameters
    ----------
    lorem ipsum dolor sit amet - lorem ipsum dolor sit amet

    Returns
    -------
    lorem ipsum dolor sit amet

    '''
    del_shape = shape._element
    slide.shapes.element.remove(del_shape)


def replace_text_by_image(ppt, search_str, image_link, how='all', adjust_size=False):
    """
    Build to change keyswords to data, in this case images
    Parameters
    ----------
    ppt : Presentation type
        slides
    search_str : Text
        key we need to search
    image_link : image path
        image path for substitution
    slide_n : int or None
        if has a number just replace on the slide with this index

    Returns
    -------
    pptx with the data in the places of the keys

    """
    assert (how in ['all', 'first'] or type(how) == int)

    # desired_h = inv_cm(4.09) #cm
    # desired_w = inv_cm(18.25) #cm
    desired_w = Cm(20)  # cm

    if type(how) == int:
        slide_loop = [ppt.slides[how]]
    else:
        slide_loop = ppt.slides

    for slide in slide_loop:
        for shape in slide.shapes:
            if shape.has_text_frame:
                if (shape.text.find(search_str)) != -1:
                    # break

                    if adjust_size:
                        size = Image.open(image_link).size
                        width, height = size
                        ratio = height / width

                        width = desired_w
                        height = width * ratio

                        left = (SLIDE_WIDTH / 2) - (width / 2)
                        top = (SLIDE_HEIGHT / 2) - (height / 2)

                    else:

                        width, height = shape.width, shape.height
                        left, top = shape.left, shape.top

                    pic = slide.shapes.add_picture(
                        image_link, left=left, top=top, height=height, width=width)
                    delete_shape(slide, shape)

                    if how == 'first':
                        return ppt

    return ppt


def duplicate_slides(ppt1, ppt2='same', index=-1, n_copies=1):
    if ppt2 == 'same' or ppt1 == ppt2:
        ppt1.save('temp.pptx')
        ppt1 = slides.Presentation('temp.pptx')
        ppt2 = ppt1
    else:
        ppt1.save('temp.pptx')
        ppt1 = slides.Presentation('temp.pptx')

        ppt2.save('temp.pptx')
        ppt2 = slides.Presentation('temp.pptx')

    for i in range(n_copies):
        ppt1.slides.add_clone(ppt2.slides[index])
        # ppt1.slides[-1].shapes[-1].text_frame.text = ''
        # Lib diferente da pptx então nosso delete_shape não funciona
        ppt1.slides[-1].shapes.remove(ppt1.slides[-1].shapes[-1])

    ppt1.save('temp.pptx', slides.export.SaveFormat.PPTX)

    ppt = Presentation('temp.pptx')

    for slide in ppt.slides:
        # slide.shapes[-1].text = ''
        delete_shape(slide, slide.shapes[-1])

    os.remove('temp.pptx')

    return ppt


def search_index(ret, ppt, test_string, how='first'):
    assert (how in ['all', 'first'] or type(how) == int)

    if type(how) == int:
        slide_loop = [ppt.slides[how]]
    else:
        slide_loop = ppt.slides

    for slide in slide_loop:
        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape.text.find(test_string) != -1:
                    text_frame = shape.text_frame

                    for para in text_frame.paragraphs:
                        try:
                            cur_text = para.runs[0].text

                            if str(test_string) in cur_text:
                                index = ppt.slides.index(slide)
                                if ret == 'index':
                                    return index
                                elif ret == 'slide':
                                    return slide
                        except:
                            pass


def duplicate_organize(ppt, key, qtd, dic):
    '''
    built to search for a slide with a specific key, duplicate it n times and change generic keys to keys in sequence.
    exemple filter_x to filter_1 in sequence of slides
    Parameters
    ----------
    ppt : Presentation type
    key : String type
    qtd : Int type
    dic : Dict type

    Returns
    -------
    new ppt and list with titles
    '''
    lista_r = []
    slide_id = search_index(ret='index', ppt=ppt, test_string=key)
    print(f"Copiando {qtd - 1} vezes o slide {slide_id + 1}")
    ppt = duplicate_slides(ppt1=ppt, ppt2='same', index=slide_id, n_copies=qtd - 1)
    for k in dic:
        key = k['name']
        sub = k['sub']
        pk = k['pk']
        lista = []

        for i in range(qtd):
            i += 1
            a = '{' + f'{sub}_{i}' + '}'
            print('a', a)
            lista.append(a)

        if pk != 1:
            lista_r += lista
        for x in lista:
            ppt = replace_text(ppt, key, x, 'first')
    return ppt, lista_r


df = pd.read_excel('unidades.xlsx')
input_ = 'etiquetas.pptx'
padrao = 11
#  Abrindo template com lib
ppt = Presentation(input_)

linhas = len(df)
copias = linhas // padrao
if linhas % padrao != 0:
    copias += 1
keys = []
words = []
print(f"Duplicando o slide {copias - 1} vezes")
ppt = duplicate_slides(ppt, n_copies=copias - 1)
print("Cópias OK")
print("Processo de lista")
for index, row in df.iterrows():
    keys.append('{nome}')
    words.append(row['nome'])
    keys.append('{unidade}')
    words.append(row['unidade'])
print("Listas OK")
print("Processo de substituição")
for search_string, replace_string in zip(keys, words):
    ppt = replace_text(ppt, search_string, replace_string, how='first')
ppt.save('etiquetas_finalizado.pptx')
print("done!")

